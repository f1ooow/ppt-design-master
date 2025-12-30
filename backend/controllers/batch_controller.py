"""
批量任务控制器
处理文件上传、任务创建、生成控制等 API
"""
import os
import io
import base64
import logging
import threading
from flask import Blueprint, request, current_app, send_file
from werkzeug.utils import secure_filename

from config import Config
from utils.response import success_response, error_response, created_response
from services.file_parser import FileParser, ScriptPage
from services.task_manager import get_task_manager, TaskStatus
from services.ai_service import get_ai_service

logger = logging.getLogger(__name__)

batch_bp = Blueprint('batch', __name__)

# 允许的文件扩展名
ALLOWED_EXTENSIONS = {'.xlsx', '.xls', '.txt', '.md', '.json'}


def allowed_file(filename: str) -> bool:
    """检查文件扩展名是否允许"""
    return os.path.splitext(filename)[1].lower() in ALLOWED_EXTENSIONS


@batch_bp.route('/upload', methods=['POST'])
def upload_script():
    """
    上传脚本文件并解析

    Returns:
        解析后的页面列表
    """
    if 'file' not in request.files:
        return error_response("没有上传文件", 400)

    file = request.files['file']
    if file.filename == '':
        return error_response("文件名为空", 400)

    if not allowed_file(file.filename):
        return error_response(f"不支持的文件格式，支持: {', '.join(ALLOWED_EXTENSIONS)}", 400)

    try:
        # 保存文件
        filename = secure_filename(file.filename)
        # 处理中文文件名
        if not filename or filename == '':
            filename = file.filename

        upload_path = os.path.join(Config.UPLOAD_FOLDER, filename)
        file.save(upload_path)
        logger.info(f"文件上传成功: {upload_path}")

        # 解析文件
        pages = FileParser.parse_file(upload_path)

        return success_response({
            'filename': filename,
            'total_pages': len(pages),
            'pages': [p.to_dict() for p in pages]
        }, "文件解析成功")

    except Exception as e:
        logger.error(f"文件处理失败: {e}")
        return error_response(f"文件处理失败: {str(e)}", 500)


@batch_bp.route('/create', methods=['POST'])
def create_task():
    """
    创建批量任务

    Request body:
        {
            "name": "任务名称",
            "pages": [ScriptPage 列表]
        }
    """
    data = request.get_json()
    if not data:
        return error_response("请求数据为空", 400)

    name = data.get('name', '未命名任务')
    pages_data = data.get('pages', [])

    if not pages_data:
        return error_response("页面列表为空", 400)

    try:
        # 转换为 ScriptPage 对象
        pages = []
        for i, p in enumerate(pages_data):
            page = ScriptPage(
                index=i,
                shot_number=p.get('shot_number', str(i + 1)),
                segment=p.get('segment', ''),
                narration=p.get('narration', ''),
                visual_hint=p.get('visual_hint', '')
            )
            pages.append(page)

        # 创建任务
        task_manager = get_task_manager()
        task = task_manager.create_task(name, pages)

        return created_response({
            'task_id': task.id,
            'name': task.name,
            'total_pages': task.total_pages
        }, "任务创建成功")

    except Exception as e:
        logger.error(f"创建任务失败: {e}")
        return error_response(f"创建任务失败: {str(e)}", 500)


@batch_bp.route('/<task_id>/generate', methods=['POST'])
def start_generation(task_id: str):
    """
    开始生成（描述 + 图片）

    Args:
        task_id: 任务 ID
    """
    task_manager = get_task_manager()
    task = task_manager.get_task(task_id)

    if not task:
        return error_response("任务不存在", 404)

    if task.status not in [TaskStatus.PENDING, TaskStatus.ERROR]:
        return error_response(f"任务状态不允许启动: {task.status.value}", 400)

    # 在后台线程中执行生成
    def run_generation():
        ai_service = get_ai_service()

        # 生成描述
        def generate_description(page: ScriptPage) -> str:
            return ai_service.generate_page_description(
                shot_number=page.shot_number,
                segment=page.segment,
                narration=page.narration,
                visual_hint=page.visual_hint
            )

        task_manager.run_descriptions_generation(task_id, generate_description)

        # 生成图片（暂时跳过，等图片生成 API 完善后启用）
        # def generate_image(page: ScriptPage) -> str:
        #     image = ai_service.generate_image(page.description)
        #     if image:
        #         image_path = os.path.join(Config.OUTPUT_FOLDER, f"{task_id}_{page.index}.png")
        #         image.save(image_path)
        #         return image_path
        #     return ""
        #
        # task_manager.run_images_generation(task_id, generate_image)

        # 直接标记为完成（描述生成完成即可）
        task_manager.update_task_status(task_id, TaskStatus.COMPLETED, "描述生成完成")

    thread = threading.Thread(target=run_generation, daemon=True)
    thread.start()

    return success_response({
        'task_id': task_id,
        'status': 'started'
    }, "生成任务已启动")


@batch_bp.route('/<task_id>/status', methods=['GET'])
def get_task_status(task_id: str):
    """
    获取任务状态

    Args:
        task_id: 任务 ID
    """
    task_manager = get_task_manager()
    task = task_manager.get_task(task_id)

    if not task:
        return error_response("任务不存在", 404)

    return success_response(task.to_dict())


@batch_bp.route('/<task_id>', methods=['DELETE'])
def delete_task(task_id: str):
    """
    删除任务

    Args:
        task_id: 任务 ID
    """
    task_manager = get_task_manager()

    if task_manager.delete_task(task_id):
        return success_response(message="任务已删除")
    else:
        return error_response("任务不存在", 404)


@batch_bp.route('/list', methods=['GET'])
def list_tasks():
    """获取所有任务列表"""
    task_manager = get_task_manager()
    tasks = task_manager.get_all_tasks()

    return success_response({
        'total': len(tasks),
        'tasks': [t.to_dict() for t in tasks]
    })


@batch_bp.route('/generate-description', methods=['POST'])
def generate_single_description():
    """
    生成单页描述

    Request body:
        {
            "narration": "讲稿内容",
            "visual_hint": "画面提示（可选）",
            "full_context": "完整脚本上下文（可选）",
            "current_index": 当前页索引（可选）,
            "custom_prompt": "自定义提示词（可选，前端传来）",
            "template_base64": "模板图片base64（可选）"
        }
    """
    data = request.get_json()
    if not data:
        return error_response("请求数据为空", 400)

    narration = data.get('narration', '')
    visual_hint = data.get('visual_hint', '')
    full_context = data.get('full_context', '')  # 完整上下文
    current_index = data.get('current_index')     # 当前页索引
    custom_prompt = data.get('custom_prompt')     # 自定义提示词
    template_base64 = data.get('template_base64') # 模板图片

    # 调试日志
    logger.info(f"[描述生成] 收到请求: narration长度={len(narration)}, 有custom_prompt={custom_prompt is not None}, 有template={template_base64 is not None and len(template_base64) > 0}")

    if not narration:
        return error_response("讲稿内容不能为空", 400)

    try:
        ai_service = get_ai_service()
        description = ai_service.generate_page_description(
            shot_number="",
            segment="",
            narration=narration,
            visual_hint=visual_hint,
            full_context=full_context if full_context else None,
            current_index=current_index,
            custom_prompt=custom_prompt,
            template_base64=template_base64
        )

        return success_response({
            'description': description
        }, "描述生成成功")

    except Exception as e:
        logger.error(f"生成描述失败: {e}")
        return error_response(f"生成描述失败: {str(e)}", 500)


@batch_bp.route('/generate-image', methods=['POST'])
def generate_single_image():
    """
    生成单页 PPT 图片

    Request body:
        {
            "narration": "讲稿内容",
            "description": "页面描述（可选）",
            "page_type": "content/cover/ending",
            "aspect_ratio": "16:9",
            "template_base64": "模板图片base64（可选）",
            "custom_prompt": "自定义提示词（可选）",
            "api_config": {  // 可选，前端传来的 API 配置
                "api_url": "...",
                "api_key": "...",
                "model": "..."
            }
        }
    """
    data = request.get_json()
    if not data:
        return error_response("请求数据为空", 400)

    narration = data.get('narration', '')
    description = data.get('description', '')
    page_type = data.get('page_type', 'content')
    aspect_ratio = data.get('aspect_ratio', '16:9')
    template_base64 = data.get('template_base64')  # 模板图片
    custom_prompt = data.get('custom_prompt')  # 自定义提示词
    api_config = data.get('api_config')  # 前端传来的 API 配置

    logger.info(f"[图片生成] 收到请求: 有custom_prompt={custom_prompt is not None}, 有template={template_base64 is not None}")

    if not narration and not description:
        return error_response("讲稿或描述至少需要一个", 400)

    try:
        ai_service = get_ai_service()

        # 如果有前端配置，临时覆盖 AI 服务的配置
        if api_config and api_config.get('api_url') and api_config.get('api_key'):
            logger.info(f"[图片API] 使用前端传来的 API 配置: {api_config.get('api_url')}")
            # 临时保存原配置
            original_api_base = ai_service.image_api_base
            original_api_key = ai_service.image_api_key
            original_model = ai_service.image_model

            # 使用前端配置
            ai_service.image_api_base = api_config.get('api_url')
            ai_service.image_api_key = api_config.get('api_key')
            if api_config.get('model'):
                ai_service.image_model = api_config.get('model')

            try:
                image_base64 = ai_service.generate_ppt_image(
                    narration=narration,
                    description=description,
                    page_type=page_type,
                    aspect_ratio=aspect_ratio,
                    template_base64=template_base64,
                    custom_prompt=custom_prompt
                )
            finally:
                # 恢复原配置
                ai_service.image_api_base = original_api_base
                ai_service.image_api_key = original_api_key
                ai_service.image_model = original_model
        else:
            image_base64 = ai_service.generate_ppt_image(
                narration=narration,
                description=description,
                page_type=page_type,
                aspect_ratio=aspect_ratio,
                template_base64=template_base64,
                custom_prompt=custom_prompt
            )

        if image_base64:
            return success_response({
                'image_base64': image_base64
            }, "图片生成成功")
        else:
            return error_response("图片生成失败，未能获取有效图片", 500)

    except Exception as e:
        logger.error(f"生成图片失败: {e}")
        return error_response(f"生成图片失败: {str(e)}", 500)


@batch_bp.route('/extract-image', methods=['POST'])
def extract_single_image():
    """
    从裁剪的图片中提取插画

    Request body:
        {
            "cropped_image_base64": "裁剪后的图片base64数据"
        }
    """
    data = request.get_json()
    if not data:
        return error_response("请求数据为空", 400)

    cropped_image_base64 = data.get('cropped_image_base64', '')

    if not cropped_image_base64:
        return error_response("裁剪图片不能为空", 400)

    try:
        ai_service = get_ai_service()
        image_base64 = ai_service.extract_illustration(cropped_image_base64)

        if image_base64:
            return success_response({
                'image': image_base64
            }, "插画提取成功")
        else:
            return error_response("插画提取失败，未能获取有效图片", 500)

    except Exception as e:
        logger.error(f"提取插画失败: {e}")
        return error_response(f"提取插画失败: {str(e)}", 500)


@batch_bp.route('/remove-background', methods=['POST'])
def remove_single_background():
    """
    去除单张图片的模板背景
    """
    data = request.get_json()
    if not data:
        return error_response("请求数据为空", 400)

    image_base64 = data.get('image_base64', '')

    if not image_base64:
        return error_response("图片数据不能为空", 400)

    try:
        ai_service = get_ai_service()
        result_base64 = ai_service.remove_template_background(image_base64)

        if result_base64:
            return success_response({
                'image_base64': result_base64
            }, "去背景成功")
        else:
            return error_response("去背景失败，未能获取有效图片", 500)

    except Exception as e:
        logger.error(f"去背景失败: {e}")
        return error_response(f"去背景失败: {str(e)}", 500)


@batch_bp.route('/clean-image', methods=['POST'])
def clean_single_image():
    """
    清洗单张PPT图片：去除文字和模板装饰，只保留插图/图标

    Request body:
        {
            "image_base64": "原始PPT图片的base64数据"
        }
    """
    data = request.get_json()
    if not data:
        return error_response("请求数据为空", 400)

    image_base64 = data.get('image_base64', '')

    if not image_base64:
        return error_response("图片数据不能为空", 400)

    try:
        ai_service = get_ai_service()
        result_base64 = ai_service.clean_slide_image(image_base64)

        if result_base64:
            return success_response({
                'image_base64': result_base64
            }, "图片清洗成功")
        else:
            return error_response("图片清洗失败，未能获取有效图片", 500)

    except Exception as e:
        logger.error(f"图片清洗失败: {e}")
        return error_response(f"图片清洗失败: {str(e)}", 500)


@batch_bp.route('/export-ppt', methods=['POST'])
def export_ppt():
    """
    导出 PPT 文件

    Request body:
        {
            "pages": [
                {
                    "type": "cover/content/ending",
                    "title": "标题",
                    "subtitle": "副标题",
                    "narration": "讲稿",
                    "image_base64": "base64图片数据"
                }
            ]
        }
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN

    logger.info("[导出PPT] 开始处理导出请求")

    data = request.get_json()
    if not data:
        logger.error("[导出PPT] 请求数据为空")
        return error_response("请求数据为空", 400)

    pages_data = data.get('pages', [])
    if not pages_data:
        logger.error("[导出PPT] 页面列表为空")
        return error_response("页面列表为空", 400)

    logger.info(f"[导出PPT] 收到 {len(pages_data)} 页待导出")

    try:
        # 创建 16:9 的 PPT
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 宽度
        prs.slide_height = Inches(7.5)    # 16:9 高度

        successful_pages = 0
        failed_pages = []

        for idx, page in enumerate(pages_data):
            page_type = page.get('type', 'content')
            image_base64 = page.get('image_base64')

            logger.info(f"[导出PPT] 处理第 {idx + 1} 页，类型: {page_type}")

            # 添加空白幻灯片
            blank_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(blank_layout)

            if image_base64:
                try:
                    # 解析 base64 图片
                    if image_base64.startswith('data:'):
                        # 移除 data:image/xxx;base64, 前缀
                        image_base64 = image_base64.split(',', 1)[1]

                    image_data = base64.b64decode(image_base64)
                    image_stream = io.BytesIO(image_data)

                    # 图片铺满整个幻灯片
                    slide.shapes.add_picture(
                        image_stream,
                        Inches(0),
                        Inches(0),
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                    successful_pages += 1
                    logger.info(f"[导出PPT] 第 {idx + 1} 页图片添加成功")
                except Exception as e:
                    logger.error(f"[导出PPT] 第 {idx + 1} 页图片处理失败: {e}")
                    failed_pages.append(idx + 1)
            else:
                logger.warning(f"[导出PPT] 第 {idx + 1} 页没有图片数据")
                failed_pages.append(idx + 1)

        logger.info(f"[导出PPT] 处理完成，成功: {successful_pages}，失败: {len(failed_pages)}")

        # 保存到内存
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        logger.info("[导出PPT] PPT文件生成成功，准备发送")

        return send_file(
            output,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='presentation.pptx'
        )

    except Exception as e:
        logger.error(f"[导出PPT] 导出失败: {e}", exc_info=True)
        return error_response(f"导出PPT失败: {str(e)}", 500)
