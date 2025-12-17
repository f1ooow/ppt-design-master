"""
PPTX 导出服务
将生成的图片导出为 PowerPoint 文件
"""
import os
import logging
from pathlib import Path
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Emu
from PIL import Image

from config import Config
from .file_parser import ScriptPage

logger = logging.getLogger(__name__)


class ExportService:
    """PPTX 导出服务"""

    # 16:9 幻灯片尺寸（单位：英寸）
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    def __init__(self):
        self.output_folder = Config.OUTPUT_FOLDER
        os.makedirs(self.output_folder, exist_ok=True)

    def export_to_pptx(self, pages: List[ScriptPage], output_name: str,
                       include_notes: bool = True) -> str:
        """
        导出为 PPTX 文件

        Args:
            pages: 页面列表
            output_name: 输出文件名（不含扩展名）
            include_notes: 是否在备注中包含讲稿内容

        Returns:
            输出文件路径
        """
        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = self.SLIDE_WIDTH
        prs.slide_height = self.SLIDE_HEIGHT

        # 使用空白布局
        blank_layout = prs.slide_layouts[6]  # 空白布局

        for page in pages:
            if not page.image_path or not os.path.exists(page.image_path):
                logger.warning(f"页面 {page.index} 没有图片，跳过")
                continue

            # 添加幻灯片
            slide = prs.slides.add_slide(blank_layout)

            # 获取图片尺寸并计算缩放
            with Image.open(page.image_path) as img:
                img_width, img_height = img.size

            # 计算缩放比例，保持宽高比并填充幻灯片
            slide_width_px = self.SLIDE_WIDTH.emu
            slide_height_px = self.SLIDE_HEIGHT.emu

            scale_w = slide_width_px / Emu(Inches(img_width / 96).emu)
            scale_h = slide_height_px / Emu(Inches(img_height / 96).emu)
            scale = max(scale_w, scale_h)

            # 计算图片在幻灯片上的位置（居中）
            final_width = Emu(int(Inches(img_width / 96).emu * scale))
            final_height = Emu(int(Inches(img_height / 96).emu * scale))

            left = Emu(int((slide_width_px - final_width) / 2))
            top = Emu(int((slide_height_px - final_height) / 2))

            # 添加图片，铺满整个幻灯片
            slide.shapes.add_picture(
                page.image_path,
                left=Emu(0),
                top=Emu(0),
                width=self.SLIDE_WIDTH,
                height=self.SLIDE_HEIGHT
            )

            # 添加备注（讲稿内容）
            if include_notes and page.narration:
                notes_slide = slide.notes_slide
                notes_frame = notes_slide.notes_text_frame
                notes_frame.text = page.narration

        # 保存文件
        output_path = os.path.join(self.output_folder, f"{output_name}.pptx")
        prs.save(output_path)
        logger.info(f"PPTX 导出成功: {output_path}")

        return output_path

    def export_images_only(self, pages: List[ScriptPage], output_name: str) -> str:
        """
        仅导出图片到文件夹

        Args:
            pages: 页面列表
            output_name: 输出文件夹名

        Returns:
            输出文件夹路径
        """
        output_dir = os.path.join(self.output_folder, output_name)
        os.makedirs(output_dir, exist_ok=True)

        for page in pages:
            if not page.image_path or not os.path.exists(page.image_path):
                continue

            # 复制图片到输出目录
            src_path = Path(page.image_path)
            dst_path = Path(output_dir) / f"slide_{page.index + 1:03d}{src_path.suffix}"

            with open(src_path, 'rb') as src, open(dst_path, 'wb') as dst:
                dst.write(src.read())

        logger.info(f"图片导出成功: {output_dir}")
        return output_dir


# 单例实例
_export_service: Optional[ExportService] = None


def get_export_service() -> ExportService:
    """获取导出服务单例"""
    global _export_service
    if _export_service is None:
        _export_service = ExportService()
    return _export_service
